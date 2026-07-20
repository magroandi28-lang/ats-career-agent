# -*- coding: utf-8 -*-
"""TUDÁSBÁZIS-ÉPÍTŐ — 1. lépés: szöveg-kinyerés és darabolás.

Andi munka- és szervezetpszichológia anyagaiból kb. 1500 karakteres,
forrás-megjelölt szakaszokat készít. Kimenet: db/tudasbazis_nyers.json
(A 2. lépés — Gemini embedding + Supabase feltöltés — külön script lesz.)

Futtatás:  python scripts/tudasbazis_epito.py "<mappa útvonala>"
"""

import json
import os
import re
import sys
import warnings

warnings.filterwarnings("ignore")

# Mely fájlok kerülnek a tudásbázisba (a mappából):
FORRASOK = [
    ("Munka és szervezetpszichológia tankönyv 2011.pdf", "Corvinus tankönyv (2011)"),
    ("Treningkonyv.pdf", "Tréneri kézikönyv (2017)"),
    ("modszertan-teljes.pdf", "Játék- és módszertani gyűjtemény"),
    ("vdocuments.mx_rudas-janos-delfi-oeroekoesei-olvasom.pdf", "Rudas: Delfi örökösei"),
    ("state-of-the-global-workplace-2022-download.pdf", "Gallup: State of the Global Workplace 2022"),
    ("1 MH Szervezetpszicho területe coo.pptx", "Dia: A szervezetpszichológia területei"),
    ("2MH Vezetéselméletek coo.pptx", "Dia: Vezetéselméletek"),
    ("3MH Munka és karrier coo.pptx", "Dia: Munka és karrier"),
    ("4MH Motiváció coo.pptx", "Dia: Motiváció"),
    ("5MH Szervezeti kultúra 2 coo (1).pptx", "Dia: Szervezeti kultúra"),
    ("6MH Pozitív szervezetpszichológia coo.pptx", "Dia: Pozitív szervezetpszichológia"),
    ("Munkapszichológia esszencia .docx", "Munkapszichológia esszencia"),
    ("Munkapszichológia álláshirdetéshez.docx", "Munkapszichológia (álláshirdetéshez)"),
    ("Órai jegyzet.docx", "Órai jegyzet"),
    ("összegzés.docx", "Összegzés"),
]

CHUNK_MERET = 1500   # karakter
MIN_MERET = 300      # ennél rövidebb szakaszt eldobunk (zaj)


def tisztit(szoveg: str) -> str:
    szoveg = re.sub(r"\(cid:\d+\)", "", szoveg)          # hibás PDF-karakterek
    szoveg = re.sub(r"[ \t]+", " ", szoveg)
    szoveg = re.sub(r"\n{3,}", "\n\n", szoveg)
    return szoveg.strip()


def pdf_szoveg(ut: str) -> str:
    try:
        import fitz  # PyMuPDF — sokkal gyorsabb
        doc = fitz.open(ut)
        return "\n".join(p.get_text() for p in doc)
    except ImportError:
        import pdfplumber
        resz = []
        with pdfplumber.open(ut) as pdf:
            for p in pdf.pages:
                resz.append(p.extract_text() or "")
        return "\n".join(resz)


def docx_szoveg(ut: str) -> str:
    from docx import Document
    d = Document(ut)
    return "\n".join(p.text for p in d.paragraphs if p.text.strip())


def pptx_szoveg(ut: str) -> str:
    from pptx import Presentation
    pr = Presentation(ut)
    resz = []
    for i, dia in enumerate(pr.slides, 1):
        diaszoveg = [sh.text_frame.text for sh in dia.shapes if sh.has_text_frame]
        t = "\n".join(x for x in diaszoveg if x.strip())
        if t.strip():
            resz.append(f"[{i}. dia] {t}")
    return "\n\n".join(resz)


def darabol(szoveg: str) -> list:
    """Bekezdés-határokon vágva kb. CHUNK_MERET-es szakaszok."""
    bekezdesek = [b.strip() for b in szoveg.split("\n") if b.strip()]
    darabok, aktualis = [], ""
    for b in bekezdesek:
        if len(aktualis) + len(b) > CHUNK_MERET and len(aktualis) >= MIN_MERET:
            darabok.append(aktualis.strip())
            aktualis = b
        else:
            aktualis += "\n" + b
    if len(aktualis.strip()) >= MIN_MERET:
        darabok.append(aktualis.strip())
    return darabok


def main():
    mappa = sys.argv[1] if len(sys.argv) > 1 else "."
    kimenet = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                           "db", "tudasbazis_nyers.json")
    osszes = []
    for fajl, forras_nev in FORRASOK:
        ut = os.path.join(mappa, fajl)
        if not os.path.exists(ut):
            print(f"  KIHAGYVA (nincs meg): {fajl}")
            continue
        try:
            if fajl.lower().endswith(".pdf"):
                szoveg = pdf_szoveg(ut)
            elif fajl.lower().endswith(".docx"):
                szoveg = docx_szoveg(ut)
            elif fajl.lower().endswith(".pptx"):
                szoveg = pptx_szoveg(ut)
            else:
                continue
            szoveg = tisztit(szoveg)
            darabok = darabol(szoveg)
            for i, d in enumerate(darabok, 1):
                osszes.append({"forras": forras_nev, "resz": i, "szoveg": d})
            print(f"  OK: {forras_nev}: {len(darabok)} szakasz ({len(szoveg)} karakter)")
        except Exception as e:
            print(f"  HIBA ({fajl}): {e}")

    with open(kimenet, "w", encoding="utf-8") as f:
        json.dump(osszes, f, ensure_ascii=False, indent=1)
    print(f"\nKESZ: {len(osszes)} szakasz -> {kimenet}")


if __name__ == "__main__":
    main()
